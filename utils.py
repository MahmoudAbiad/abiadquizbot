"""
Utility functions for file processing and hashing.
"""

import os
import hashlib
from logger import get_logger, log_error, log_info

logger = get_logger(__name__)

def calculate_file_hash(file_path: str) -> str:
    """
    Return a SHA-256 digest of file bytes only. File names are deliberately ignored.
    """
    hasher = hashlib.sha256()
    try:
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hasher.update(chunk)
        return hasher.hexdigest()
    except Exception as e:
        log_error(logger, f"Error calculating hash for {file_path}: {e}")
        return ""

def safe_file_cleanup(file_path: str) -> bool:
    """Safely delete a file with error handling."""
    try:
        if os.path.exists(file_path):
            os.remove(file_path)
            log_info(logger, f"File deleted: {file_path}")
            return True
        return False
    except Exception as e:
        log_error(logger, f"Error deleting file {file_path}: {e}", exception=e)
        return False

def ensure_directory_exists(directory_path: str) -> bool:
    """Ensure a directory exists, create if necessary."""
    try:
        if not os.path.exists(directory_path):
            os.makedirs(directory_path)
        return True
    except Exception as e:
        log_error(logger, f"Error creating directory {directory_path}: {e}", exception=e)
        return False

def format_file_size(size_bytes: int) -> str:
    """Format bytes to human-readable file size."""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.2f} {unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.2f} TB"

def extract_text_from_file(file_path: str) -> str | None:
    """
    استخراج النص الصافي من ملفات Word (.docx), PowerPoint (.pptx), TXT.
    ملاحظة: الصيغ القديمة (.doc, .ppt) لا تُقرأ بواسطة python-docx/pptx مباشرة.
    """
    ext = os.path.splitext(file_path)[1].lower()
    try:
        # 1. ملفات النص المباشر
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        # 2. مستندات Word الحديثة فقط
        elif ext == ".docx":
            import docx
            doc = docx.Document(file_path)
            full_text = []

            # استخراج النصوص من الفقرات
            for p in doc.paragraphs:
                if p.text.strip():
                    full_text.append(p.text.strip())

            # استخراج النصوص من الجداول
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            full_text.append(cell_text)

            return "\n".join(full_text) if full_text else None

        # 3. عروض PowerPoint الحديثة فقط
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            text_runs = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                text_runs.append(text)

            return "\n".join(text_runs) if text_runs else None

        # 4. الصيغ الثنائية القديمة (Log Warning بدلاً من العطل)
        elif ext in [".doc", ".ppt"]:
            logger.warning(f"File {file_path} is legacy binary format ({ext}) and cannot be parsed directly.")
            return None

    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        return None

    return None